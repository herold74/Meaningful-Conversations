#!/usr/bin/env python3
"""
Meaningful Conversations — PPTX Generator
==========================================

Generates a 5-slide PowerPoint presentation from the User Access Matrix.

Slides:
  1. Title
  2. Feature Access Matrix — Platform Features & Functions
  3. Feature Access Matrix — Coaching Bots & Exclusive Features
  4. Pricing Model — 4 Tiers
  5. Bot Categories — Bronze · Silver · Gold
  6. Upgrade Paths & Discounts

Usage:
  pip install python-pptx
  python3 DOCUMENTATION/generate_pptx.py

Output:
  DOCUMENTATION/Meaningful-Conversations-Access-Matrix-Pricing.pptx

Last updated: February 2026 — v1.8.9+
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Brand Colors ──────────────────────────────────────────────────────────────

TEAL = RGBColor(0x1B, 0x72, 0x72)
TEAL_LIGHT = RGBColor(0xE0, 0xF2, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
BRONZE = RGBColor(0xCD, 0x7F, 0x32)
BRONZE_BG = RGBColor(0xFD, 0xF0, 0xE0)
SILVER = RGBColor(0x47, 0x56, 0x69)          # Darker silver text for better contrast
SILVER_BG = RGBColor(0xDB, 0xE2, 0xEF)       # Blue-tinted background
GOLD = RGBColor(0xD9, 0x77, 0x06)
GOLD_BG = RGBColor(0xFF, 0xFB, 0xEB)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED = RGBColor(0xDC, 0x26, 0x26)

# ─── Helpers ───────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text,
                 font_size=12, bold=False, color=DARK,
                 alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_table_cell(table, row, col, text,
                   font_size=9, bold=False, color=DARK,
                   alignment=PP_ALIGN.CENTER, bg_color=None):
    cell = table.cell(row, col)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
    cell.margin_left = Emu(45720)
    cell.margin_right = Emu(45720)
    cell.margin_top = Emu(27432)
    cell.margin_bottom = Emu(27432)


def add_rounded_rect(slide, left, top, width, height,
                     fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.1
    return shape


def build_matrix_slide(slide, title, subtitle, rows_data):
    """Build a matrix slide with a teal header bar and a data table."""
    add_bg(slide, WHITE)
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = TEAL
    hdr.line.fill.background()
    add_text_box(slide, 0.5, 0.12, 10, 0.35, title,
                 font_size=26, bold=True, color=WHITE)
    add_text_box(slide, 0.5, 0.5, 10, 0.3, subtitle,
                 font_size=12, color=RGBColor(0xB0, 0xD8, 0xD8))

    num_rows = len(rows_data) + 1
    num_cols = 5

    table_h = min(6.2, 0.38 * num_rows + 0.1)
    tbl = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.4), Inches(1.15), Inches(12.5), Inches(table_h)).table

    tbl.columns[0].width = Inches(4.8)
    tbl.columns[1].width = Inches(1.925)
    tbl.columns[2].width = Inches(1.925)
    tbl.columns[3].width = Inches(1.925)
    tbl.columns[4].width = Inches(1.925)

    headers = [
        'Feature',
        'Guest\n(kostenlos)',
        'Registered\n(3,90 €/Monat)',
        'Premium\n(9,90 €/Monat)',
        'Client\n(durch Coach)',
    ]
    for ci, h in enumerate(headers):
        add_table_cell(tbl, 0, ci, h, font_size=10, bold=True,
                       color=WHITE, bg_color=TEAL)

    for ri, (label, guest, reg, prem, client, is_header, section_color) \
            in enumerate(rows_data):
        row_idx = ri + 1
        if is_header:
            bg = (GRAY_LIGHT if not section_color
                  else (BRONZE_BG if section_color == BRONZE
                        else SILVER_BG if section_color == SILVER
                        else GOLD_BG))
            fc = DARK if not section_color else section_color
            for ci in range(num_cols):
                add_table_cell(tbl, row_idx, ci, '', bg_color=bg)
            add_table_cell(tbl, row_idx, 0, label, font_size=10,
                           bold=True, color=fc, alignment=PP_ALIGN.LEFT,
                           bg_color=bg)
        else:
            row_bg = WHITE if ri % 2 == 0 else GRAY_LIGHT
            add_table_cell(tbl, row_idx, 0, label, font_size=10,
                           color=DARK, alignment=PP_ALIGN.LEFT, bg_color=row_bg)
            for ci, val in enumerate([guest, reg, prem, client], 1):
                vc = (GREEN if val == '✅'
                      else (RGBColor(0xBB, 0xBB, 0xBB) if val == '—'
                            else DARK))
                add_table_cell(tbl, row_idx, ci, val, font_size=10,
                               color=vc, bg_color=row_bg)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════

slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1, WHITE)

header = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
header.fill.solid()
header.fill.fore_color.rgb = TEAL
header.line.fill.background()

add_text_box(slide1, 1.5, 2.0, 10, 1.2,
             'Meaningful Conversations',
             font_size=42, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, 1.5, 3.2, 10, 0.8,
             'User Access Matrix & Preismodell',
             font_size=24, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, 1.5, 4.3, 10, 0.5,
             'by manualmode.at  •  v1.8.9  •  Februar 2026',
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

bottom = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(7.35), prs.slide_width, Inches(0.15))
bottom.fill.solid()
bottom.fill.fore_color.rgb = TEAL
bottom.line.fill.background()


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PLATFORM FEATURES & FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════════

slide2 = prs.slides.add_slide(prs.slide_layouts[6])

#                  (label, guest, registered, premium, client, is_header, section_color)
features_data = [
    ('Core Functions',                       '', '', '', '', True,  None),
    ('Chat & Voice (Web Speech API)',        '✅', '✅', '✅', '✅', False, None),
    ('Server TTS (High Quality)',            '—',  '✅', '✅', '✅', False, None),
    ('Life Context',                         'Lokal', 'Cloud E2EE', 'Cloud E2EE', 'Cloud E2EE', False, None),
    ('Cloud-Sync & Geräteübergreifend',      '—',  '✅', '✅', '✅', False, None),
    ('Persönlichkeitsprofil (OCEAN)',         '—',  '✅', '✅', '✅', False, None),
    ('Persönlichkeitsprofil (Riemann & SD)', '—',  '—',  '✅', '✅', False, None),
    ('Narrative Signature & PDF-Export',      '—',  '✅', '✅', '✅', False, None),
    ('DPC (Dynamic Prompt Composition)',      '—',  '✅', '✅', '✅', False, None),
    ('DPFL (Adaptive Learning)',              '—',  '—',  '✅', '✅', False, None),
    ('Comfort Check',                         '—',  '—',  '✅', '✅', False, None),
    ('Gamification (XP, Levels)',             'Lokal', '✅', '✅', '✅', False, None),
    ('Kalenderexport (.ics)',                 '✅', '✅', '✅', '✅', False, None),
    ('Krisenreaktion (Helplines)',            '✅', '✅', '✅', '✅', False, None),
]

build_matrix_slide(slide2,
                   'Feature Access Matrix',
                   'Plattform-Features & Funktionen',
                   features_data)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — COACHING BOTS & EXCLUSIVE FEATURES
# ═══════════════════════════════════════════════════════════════════════════════

slide3 = prs.slides.add_slide(prs.slide_layouts[6])

bots_data = [
    ('Management & Kommunikation',           '', '', '', '', True, BRONZE),
    ('Nobody (GPS, Problemlösung)',          '✅', '✅', '✅', '✅', False, None),
    ('Gloria Interview (Strukturierte Interviews)', '—', '✅', '✅', '✅', False, None),
    ('Gloria (Onboarding)',                  '✅', '✅', '✅', '✅', False, None),
    ('Coaching Bots',                        '', '', '', '', True, SILVER),
    ('Max (Ambitioniert)',                   '✅', '✅', '✅', '✅', False, None),
    ('Ava (Strategisch)',                    '✅', '✅', '✅', '✅', False, None),
    ('Kenji (Stoisch)',                      '—',  '—',  '✅', '✅', False, None),
    ('Chloe (Strukturierte Reflexion)',      '—',  '—',  '✅', '✅', False, None),
    ('Exklusiv für Klienten',               '', '', '', '', True, GOLD),
    ('Rob (Mentale Fitness)',                '—',  '—',  '—',  '✅', False, None),
    ('Victor (Systemisch)',                  '—',  '—',  '—',  '✅', False, None),
    ('Premium+ Features',                    '', '', '', '', True, None),
    ('Transcript Evaluation & PDF',          '—',  '—',  '✅', '✅', False, None),
    ('Bot-Empfehlungen (in Evaluation)',     '—',  '—',  '✅', '✅', False, None),
    ('PEP Lösungsblockaden (Dr. Bohne)',     '—',  '—',  '—',  '✅', False, None),
]

build_matrix_slide(slide3,
                   'Feature Access Matrix',
                   'Coaching-Bots & exklusive Features',
                   bots_data)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PRICING MODEL
# ═══════════════════════════════════════════════════════════════════════════════

slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, WHITE)

hdr4 = slide4.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
hdr4.fill.solid()
hdr4.fill.fore_color.rgb = TEAL
hdr4.line.fill.background()
add_text_box(slide4, 0.5, 0.15, 12, 0.6,
             'Preismodell — 4 Stufen',
             font_size=28, bold=True, color=WHITE)

tiers = [
    {
        'name': 'Guest',
        'price': 'Kostenlos',
        'color': GRAY,
        'bg': GRAY_LIGHT,
        'features': [
            'Nobody, Gloria, Max, Ava',
            'Chat & Voice',
            'Life Context (lokal)',
            'Kalenderexport',
        ],
        'limit': 'Daten nur im Browser',
    },
    {
        'name': 'Registered',
        'price': '3,90 €/Monat',
        'price2': 'oder 14,90 € einmalig',
        'color': TEAL,
        'bg': TEAL_LIGHT,
        'features': [
            'Alles aus Guest, plus:',
            'Gloria Interview (Transkript)',
            'Cloud-Sync & E2EE',
            'Server-TTS (High Quality)',
            'OCEAN-Profil, Signature, DPC',
        ],
        'limit': 'Kein Riemann/SD, DPFL, Kenji/Chloe',
    },
    {
        'name': 'Premium',
        'price': '9,90 €/Monat',
        'price2': '24,90 €/3M · 79,90 €/Jahr',
        'color': SILVER,
        'bg': RGBColor(0xE8, 0xEC, 0xF4),
        'features': [
            'Alles aus Registered, plus:',
            'Kenji & Chloe',
            'Riemann-Thomann & Spiral Dynamics',
            'DPFL & Comfort Check',
            'Transcript Evaluation (PDF & Bots)',
        ],
        'limit': 'Kein Rob/Victor, kein PEP',
    },
    {
        'name': 'Client',
        'price': 'Durch Coach',
        'price2': 'Nicht käuflich',
        'color': GOLD,
        'bg': GOLD_BG,
        'features': [
            'Alles aus Premium, plus:',
            'Rob & Victor',
            'PEP Lösungsblockaden',
            '',
            '',
        ],
        'limit': 'Vollzugang',
    },
]

card_w = 2.85
card_h = 4.2
gap = 0.25
total_w = 4 * card_w + 3 * gap
start_x = (13.333 - total_w) / 2

for i, tier in enumerate(tiers):
    x = start_x + i * (card_w + gap)
    y = 1.3

    add_rounded_rect(slide4, x, y, card_w, card_h, tier['bg'], tier['color'])

    add_text_box(slide4, x + 0.15, y + 0.15, card_w - 0.3, 0.4,
                 tier['name'],
                 font_size=20, bold=True, color=tier['color'],
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide4, x + 0.15, y + 0.55, card_w - 0.3, 0.35,
                 tier['price'],
                 font_size=16, bold=True, color=DARK,
                 alignment=PP_ALIGN.CENTER)

    if 'price2' in tier:
        add_text_box(slide4, x + 0.15, y + 0.85, card_w - 0.3, 0.25,
                     tier['price2'],
                     font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)

    div_y = y + 1.15
    div = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x + 0.3), Inches(div_y), Inches(card_w - 0.6), Pt(1))
    div.fill.solid()
    div.fill.fore_color.rgb = tier['color']
    div.line.fill.background()

    feat_y = div_y + 0.15
    for fi, feat in enumerate(tier['features']):
        if not feat:
            continue
        prefix = '✓ ' if fi > 0 or i == 0 else ''
        fc = DARK if fi > 0 or i == 0 else tier['color']
        fb = fi == 0 and i > 0
        add_text_box(slide4, x + 0.2, feat_y + fi * 0.33, card_w - 0.4, 0.3,
                     prefix + feat, font_size=10, bold=fb, color=fc)

    add_text_box(slide4, x + 0.15, y + card_h - 0.45, card_w - 0.3, 0.35,
                 tier['limit'],
                 font_size=8, color=GRAY, alignment=PP_ALIGN.CENTER)

# Upgrade arrows
arrow_y = 1.3 + card_h + 0.3
for i in range(3):
    ax = start_x + (i + 0.5) * (card_w + gap) + card_w * 0.5 - gap * 0.5
    arrow = slide4.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(ax - 0.2), Inches(arrow_y), Inches(0.5), Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = TEAL
    arrow.line.fill.background()

add_text_box(slide4, 0.5, arrow_y + 0.35, 12.3, 0.4,
             'Guest → Registered → Premium → Client  |  '
             'Natürlicher Upgrade-Pfad durch erlebten Mehrwert',
             font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — BOT CATEGORIES
# ═══════════════════════════════════════════════════════════════════════════════

slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, WHITE)

hdr5 = slide5.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
hdr5.fill.solid()
hdr5.fill.fore_color.rgb = TEAL
hdr5.line.fill.background()
add_text_box(slide5, 0.5, 0.15, 12, 0.6,
             'Bot-Kategorien — Bronze · Silver · Gold',
             font_size=28, bold=True, color=WHITE)

sections = [
    {
        'title': 'Management & Kommunikation',
        'subtitle': 'Nobody & Gloria: Guest  |  Gloria Interview: Registered',
        'color': BRONZE,
        'bg': BRONZE_BG,
        'bots': [
            ('Nobody', 'GPS-Ansatz, Problemlösung, Kommunikationsanalyse'),
            ('Gloria Interview', 'Strukturierte Interviews mit Transkript-Export'),
            ('Gloria', 'Onboarding & Erstgespräch'),
        ],
    },
    {
        'title': 'Coaching',
        'subtitle': 'Max & Ava: Guest  |  Kenji & Chloe: Premium',
        'color': SILVER,
        'bg': RGBColor(0xE8, 0xEC, 0xF4),
        'bots': [
            ('Max', 'Ambitioniert, motivierend, neugierig'),
            ('Ava', 'Strategisch, entscheidend, organisiert'),
            ('Kenji 🔒', 'Stoisch, philosophisch, weise'),
            ('Chloe 🔒', 'Strukturierte Reflexion, CBT-basiert'),
        ],
    },
    {
        'title': 'Exklusiv für Klienten',
        'subtitle': 'Nur mit manualmode.at Coaching-Beziehung',
        'color': GOLD,
        'bg': GOLD_BG,
        'bots': [
            ('Rob 🔒', 'Mentale Fitness, empathisch, achtsam'),
            ('Victor 🔒', 'Systemisch, analytisch, neutral'),
        ],
    },
]

sec_w = 3.8
sec_gap = 0.35
total_sec_w = 3 * sec_w + 2 * sec_gap
sec_start_x = (13.333 - total_sec_w) / 2

for i, sec in enumerate(sections):
    sx = sec_start_x + i * (sec_w + sec_gap)
    sy = 1.3
    sec_h = 5.3

    add_rounded_rect(slide5, sx, sy, sec_w, sec_h, sec['bg'], sec['color'])

    add_text_box(slide5, sx + 0.15, sy + 0.2, sec_w - 0.3, 0.4,
                 sec['title'],
                 font_size=16, bold=True, color=sec['color'],
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide5, sx + 0.15, sy + 0.6, sec_w - 0.3, 0.3,
                 sec['subtitle'],
                 font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)

    div = slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(sx + 0.3), Inches(sy + 1.0), Inches(sec_w - 0.6), Pt(1))
    div.fill.solid()
    div.fill.fore_color.rgb = sec['color']
    div.line.fill.background()

    for bi, (name, desc) in enumerate(sec['bots']):
        by = sy + 1.2 + bi * 0.95
        add_text_box(slide5, sx + 0.25, by, sec_w - 0.5, 0.3,
                     name, font_size=14, bold=True, color=DARK)
        add_text_box(slide5, sx + 0.25, by + 0.32, sec_w - 0.5, 0.5,
                     desc, font_size=10, color=GRAY)

add_text_box(slide5, 0.5, 6.8, 12.3, 0.4,
             '🔒 = Erfordert höheren Zugang  |  '
             'Einzelne Premium-Bots können für 4,90 € permanent '
             'freigeschaltet werden',
             font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — UPGRADE PATHS & DISCOUNTS
# ═══════════════════════════════════════════════════════════════════════════════

slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, WHITE)

hdr6 = slide6.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
hdr6.fill.solid()
hdr6.fill.fore_color.rgb = TEAL
hdr6.line.fill.background()
add_text_box(slide6, 0.5, 0.12, 10, 0.35,
             'Upgrade-Pfade & Rabatte',
             font_size=26, bold=True, color=WHITE)
add_text_box(slide6, 0.5, 0.5, 10, 0.3,
             'Frühere Investitionen werden immer anerkannt — kein Buyer\'s Remorse',
             font_size=12, color=RGBColor(0xB0, 0xD8, 0xD8))

# --- Main upgrade table ---
upgrade_rows = [
    ('Ausgangslage',          '1-Monats-Pass', '3-Monats-Pass', '1-Jahres-Pass', 'Rabatt-Modell'),
    ('Normalpreis',           '9,90 €',        '24,90 €',       '79,90 €',       '—'),
    ('Registered Monatsabo',  '9,90 €',        '24,90 €',       '79,90 €',       'Pro-rata Restmonat als Guthaben'),
    ('Registered Lifetime',   '7,90 €',        '18,90 €',       '59,90 €',       '~20–25% Loyalty-Rabatt'),
    ('1 Bot-Unlock (4,90 €)', '5,00 €',        '20,00 €',       '75,00 €',       '4,90 € Anrechnung'),
    ('2 Bot-Unlocks (9,80 €)','0,10 €',        '15,10 €',       '70,10 €',       '9,80 € Anrechnung'),
    ('Lifetime + 1 Bot',      '3,00 €',        '14,00 €',       '55,00 €',       'Loyalty + Bot kumuliert'),
    ('Lifetime + 2 Bots',     '0,10 €',        '9,10 €',        '50,10 €',       'Loyalty + Bots kumuliert'),
    ('Guest → Premium',       '9,90 €',        '24,90 €',       '79,90 €',       'Kein Rabatt (enthält Registered)'),
]

num_rows = len(upgrade_rows)
num_cols = 5
tbl6 = slide6.shapes.add_table(
    num_rows, num_cols,
    Inches(0.4), Inches(1.15), Inches(12.5), Inches(5.5)).table

tbl6.columns[0].width = Inches(3.0)
tbl6.columns[1].width = Inches(2.0)
tbl6.columns[2].width = Inches(2.0)
tbl6.columns[3].width = Inches(2.0)
tbl6.columns[4].width = Inches(3.5)

# Header row
for ci, h in enumerate(upgrade_rows[0]):
    add_table_cell(tbl6, 0, ci, h, font_size=10, bold=True,
                   color=WHITE, bg_color=TEAL)

# Normalpreis row (reference)
for ci, val in enumerate(upgrade_rows[1]):
    add_table_cell(tbl6, 1, ci, val, font_size=10,
                   bold=(ci == 0),
                   color=GRAY,
                   alignment=PP_ALIGN.LEFT if ci in (0, 4) else PP_ALIGN.CENTER,
                   bg_color=RGBColor(0xF9, 0xFA, 0xFB))

# Data rows
HIGHLIGHT = RGBColor(0x16, 0x65, 0x34)  # Dark green for discounted prices
for ri in range(2, num_rows):
    row_bg = WHITE if ri % 2 == 0 else GRAY_LIGHT
    row = upgrade_rows[ri]
    # Label column
    add_table_cell(tbl6, ri, 0, row[0], font_size=10,
                   bold=True, color=DARK, alignment=PP_ALIGN.LEFT,
                   bg_color=row_bg)
    # Price columns (1-3): highlight if discounted
    for ci in range(1, 4):
        is_discounted = row[ci] != upgrade_rows[1][ci]
        add_table_cell(tbl6, ri, ci, row[ci], font_size=10,
                       bold=is_discounted,
                       color=HIGHLIGHT if is_discounted else DARK,
                       bg_color=row_bg)
    # Rabatt-Modell column
    add_table_cell(tbl6, ri, 4, row[4], font_size=9,
                   color=GRAY, alignment=PP_ALIGN.LEFT,
                   bg_color=row_bg)

# --- Principles footer ---
principles = [
    '✓ Fallback-Sicherheit: Registered Lifetime bleibt aktiv wenn Premium abläuft',
    '✓ Kumulierbar: Loyalty-Rabatt und Bot-Guthaben stapeln sich',
    '✓ Technisch umsetzbar: PayPal Custom IDs erlauben Tracking früherer Käufe',
]
for pi, p_text in enumerate(principles):
    add_text_box(slide6, 0.5, 6.55 + pi * 0.28, 12.3, 0.25,
                 p_text, font_size=9, color=GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════

script_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(
    script_dir,
    'Meaningful-Conversations-Access-Matrix-Pricing.pptx')
prs.save(output_path)
print(f'✅ PPTX saved to: {output_path}')
print(f'   6 slides: Title, Features, Bots & Exclusive, Pricing, '
      f'Bot Categories, Upgrade Paths')
